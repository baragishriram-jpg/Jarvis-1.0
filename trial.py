import speech_recognition as sr
import pyttsx3
import time
import datetime
import wikipedia
import os
import cv2
import numpy as np
import sys
import urllib.parse
import webbrowser
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import requests
from colorama import Fore,Style
from pptx.enum.text import PP_ALIGN
from reportlab.pdfgen import canvas
import win32com.client
from email.message import EmailMessage
import smtplib
from ctypes import cast, POINTER
from comtypes import CLSCTX_ALL
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
import screen_brightness_control as sbc
from openai import OpenAI
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.chrome.service import Service
import pyautogui

# --- Configuration ---
WAKE_WORD = "jarvis"
LISTENING_DURATION = 10 # seconds 
LISTENING_TIMEOUT = 2 # seconds to wait for voice input
BASE_FILE_DIR = "C:/Users/admin/Desktop/Jarvis 3.0/Database"
BASE_PPT_DIR = "C:/Users/admin/Desktop/Jarvis 3.0/Database"
EMAIL_ADDRESS = "baragishriram@gmail.com"
EMAIL_PASSWORD = "9964324995"
OPENAI_API_KEY = "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"
client = OpenAI(api_key="OPENAI_API_KEY")
FACE_FILE = "authorized_face.npy"
THRESHOLD = 45  # lower = stricter
presentation = None
ppt_path = None

def print_loop():
    while True:
        print(Fore.BLUE +"listening...",end="",flush=True)
        print(Style.RESET_ALL,end="",flush=True)
        print("",end="",flush=True)

face_cascade = cv2.CascadeClassifier(
    cv2.data.haarcascades + "haarcascade_frontalface_default.xml"
)

def capture_face():
    cam = cv2.VideoCapture(0)
    print("📷 Look at the camera...")

    while True:
        ret, frame = cam.read()
        if not ret:
            continue

        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        faces = face_cascade.detectMultiScale(gray, 1.3, 5)

        for (x, y, w, h) in faces:
            face = gray[y:y+h, x:x+w]
            face = cv2.resize(face, (200, 200))
            cam.release()
            cv2.destroyAllWindows()
            return face

        cv2.imshow("FACE LOCK", frame)
        if cv2.waitKey(1) & 0xFF == 27:
            cam.release()
            cv2.destroyAllWindows()
            sys.exit()

def register_face():
    face = capture_face()
    np.save(FACE_FILE, face)
    print("✅ Face registered successfully")

def verify_face():
    saved_face = np.load(FACE_FILE)
    face = capture_face()

    diff = cv2.absdiff(saved_face, face)
    score = np.mean(diff)
    print("🔍 Difference score:", score)

    return score < THRESHOLD

# ================= FACE LOCK =================
if not os.path.exists(FACE_FILE):
    print("🔐 No registered face found")
    register_face()
    print("🔁 Restart program")
    sys.exit()

print("🔓 Verifying face...")
if not verify_face():
    print("❌ ACCESS DENIED")
    sys.exit()

print("✅ ACCESS GRANTED")
# ============================================


# 🔥 🔓 YOUR PROTECTED CODE STARTS HERE 🔓 🔥

print("Jarvis Secure System Activated")

for i in range(5):
    print(f"Secure task running {i+1}")

# -------- CONTACTS --------
CONTACTS = {
    "vinayak": "9353202838",
    "mom": "9964324995",
    "dad": "8722511691",
    "mohit": "7338567989",
    "basavaraj": "7204859782",
    "vinayak m c": "8618975789",
    "srujan": "7204551841",
    "aunty": "9844161073",
    "sahit": "8431482327",
    "suraj": "9945674946",
    "shishir": "6362691358"
}

# ------------------- Memory ---------------------
LOG_FILE = "conversation_history.txt"

def conversation_history(sender, message):
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with open(LOG_FILE, "a", encoding="utf-8") as file:
        file.write(f"[{timestamp}] {sender}: {message}\n")

# ----------------- AI Brain --------------------
def chatgpt_search(question):
    answer = client.responses.create(
        model="gpt-4.1-mini",
        input=question
    )
    return answer.output_text


# ---------------- TIME-BASED WISH ----------------
def time_wish():
    hour = datetime.datetime.now().hour
    if 5 <= hour < 12:
        speak("Good morning BOSS.")
    elif 12 <= hour < 17:
        speak("Good afternoon BOSS.")
    elif 17 <= hour < 24:
        speak("Good evening BOSS.")
    else:
        speak("Good night BOSS.")

# --- Text-to-Speech Function ---
def speak(text, gender="female"):
    
    engine = pyttsx3.init()
    voices = engine.getProperty('voices')
    
    engine.setProperty('voice', voices[1].id)  # Female voice
    engine.say(text)
    engine.runAndWait()


# --- Speech-to-Text Function ---
def get_audio():
    
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        # Calibrate for ambient noise.
        recognizer.adjust_for_ambient_noise(source)
        try:
            print("Listening...")
            audio = recognizer.listen(source, timeout=LISTENING_TIMEOUT)
            print("Recognizing...")
            text = recognizer.recognize_google(audio)
            return text.lower()
        except sr.UnknownValueError:
            return get_audio()
        except sr.WaitTimeoutError:
            return get_audio()
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")
            return None

# ---------------- OS CONTROL ----------------
def os_control(command):
    if "open chrome" in command:
        speak("Opening Google Chrome.")
        os.startfile("C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe")

    elif "open edge" in command:
        speak("Opening Microsoft Edge.")
        os.startfile("C:\\Program Files (x86)\\Microsoft\\Edge\\Application\\msedge.exe")

    elif "open browser" in command:
        speak("Opening default browser.")
        webbrowser.open("https://www.google.com")

    elif "open youtube" in command:
        speak("Opening YouTube.")
        webbrowser.open("https://www.youtube.com")

    elif 'open google' in user_command:
        speak("Opening Google.")
        webbrowser.open("http://www.google.com")

    elif 'open instagram' in user_command:
        speak("Opening Instagram.")
        webbrowser.open("http://www.instagram.com")

    elif 'open github' in user_command:
        speak("Opening GitHub.")
        webbrowser.open("http://www.github.com")

    else:
        speak("OS command not recognized.")

# ---------------- BROWSER SEARCH ----------------
def browser_search(command):
    keywords = command.replace("search", "").replace("find", "").strip()
    if keywords:
        speak(f"Searching for {keywords}")
        query = urllib.parse.quote(keywords)
        webbrowser.open(f"https://www.google.com/search?q={query}")
    else:
        speak("What should I search for?")

# ------------- listener ---------------
def take_command(command):
    r = sr.Recognizer()
    with sr.Microphone() as source:
        r.adjust_for_ambient_noise(source, duration=2)
        print(command)
        speak(command)
        audio = r.listen(source)

    try:
        text = r.recognize_google(audio)
        print("You said:", text)
        return text.lower()
    except sr.UnknownValueError:
        speak("I didn't understand. Please say again.")
        return take_command(command)

# ---------- File Operations ----------
def create_file():
    filename = take_command("Tell the file name").replace(" ", "_")
    extension = take_command("Tell the file extension like txt, docx, pdf, python").replace(" ", "_")

    file_path = os.path.join(BASE_FILE_DIR, f"{filename}.{extension}")

    try:
        with open(file_path, "w") as f:
            f.write("File created by Jarvis voice command.")
        speak(f"{filename} dot {extension} created successfully")
    except Exception as e:
        speak("Unable to create the file")
        print(e)

def open_file():
    filename = take_command("Tell the file name").replace(" ", "_")
    extension = take_command("Tell the file extension").replace(" ", "_")

    file_path = os.path.join(BASE_FILE_DIR, f"{filename}.{extension}")

    if os.path.exists(file_path):
        os.startfile(file_path)
        speak(f"Opening {filename}")
    else:
        speak("File not found")

# ---------------- IMAGE FROM WEB ----------------
def get_image_from_web(topic):
    url = f"https://source.unsplash.com/800x600/?{topic}"
    response = requests.get(url, timeout=10)
    return BytesIO(response.content)

# ---------------- GET CONTENT FROM WEB ----------------
def get_content_from_web(topic):
    try:
        summary = wikipedia.summary(topic, sentences=4)
        return summary
    except:
        return f"Information about {topic} could not be fetched."

# ---------------- PPT CREATOR ----------------
def ppt_creator():
    global presentation, ppt_path

    topic = take_command("Please tell the presentation topic").lower()

    presentation = Presentation()
    file_name = f"{topic.replace(' ', '_')}.pptx"
    ppt_path = os.path.join(BASE_PPT_DIR, file_name)

    # TITLE SLIDE
    slide = presentation.slides.add_slide(presentation.slide_layouts[3])
    slide.shapes.title.text = topic.title()
    slide.placeholders[1].text = get_content_from_web(topic)
    slide.placeholders[1].pic = get_image_from_web(topic)

    speak(f"Presentation on {topic} created.")

    decision = take_command("Would you like to save this").lower().replace(" ","_")

    if "add_slide" in decision:
        add_slide(topic)
    elif "save_it" in decision:
        finish_presentation()

# ---------------- ADD SLIDE WITH AUTO CONTENT ----------------
def add_slide(topic):
    content_topic = take_command("Tell the slide topic").lower()

    slide = presentation.slides.add_slide(presentation.slide_layouts[3])
    slide.shapes.title.text = content_topic.title()

    slide.placeholders[1].pic = get_image_from_web(content_topic)
    slide.placeholders[1].text = get_content_from_web(content_topic)

    try:
        image = get_image_from_web(content_topic)
        slide.shapes.add_picture(
            image, Inches(4.5), Inches(1.5), width=Inches(4)
        )
    except:
        pass

    speak(f"Slide added for {content_topic}")

    decision = take_command("Do you want another slide or save").lower().replace(" ","_")

    if "add_another_slide" in decision:
        add_slide(topic)
    elif "save_it" in decision:
        finish_presentation()

# ---------------- SAVE PRESENTATION ----------------
def finish_presentation():
    global presentation, ppt_path

    if presentation is None:
        speak("Presentation is not created yet")
        return

    presentation.save(ppt_path)
    speak("Presentation saved successfully")

# ---------------- pdf converter -----------------
def convert_to_pdf():
    filename = take_command("Tell the file name").replace(" ", "_")
    ext = take_command("Tell the file extension").lower()

    input_path = os.path.join(BASE_FILE_DIR, f"{filename}.{ext}")
    output_path = os.path.join(BASE_FILE_DIR, f"{filename}.pdf")

    if not os.path.exists(input_path):
        speak("File not found")
        return

    # ---------- TXT TO PDF ----------
    if ext == "txt":
        try:
            c = canvas.Canvas(output_path)
            with open(input_path, "r", encoding="utf-8") as f:
                text = c.beginText(40, 800)
                for line in f:
                    text.textLine(line.strip())
                c.drawText(text)
            c.save()
            speak("Text file converted to PDF successfully")
        except:
            speak("Failed to convert text file")

    # ---------- WORD TO PDF ----------
    elif ext == "docx":
        try:
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(input_path)
            doc.SaveAs(output_path, FileFormat=17)
            doc.Close()
            word.Quit()
            speak("Word file converted to PDF successfully")
        except:
            speak("Word conversion failed")

    # ---------- PPTX TO PDF ----------
    elif ext == "pptx":
        try:
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
            presentation.SaveAs(output_path, 32)
            presentation.Close()
            powerpoint.Quit()
            speak("PowerPoint converted to PDF successfully")
        except:
            speak("PowerPoint conversion failed")

    else:
        speak("Unsupported file format")

# ---------------- FILE SEARCH ----------------
def find_file(filename):
    for root, _, files in os.walk(BASE_FILE_DIR):
        for file in files:
            if filename in file.lower():
                return os.path.join(root, file)
    return None


# ---------------- EMAIL SENDER ----------------
def email_sender():
    speak("Please tell the receiver email address.")
    receiver = get_audio().replace(" ", "")
    
    speak("What should I say in the email?")
    message_text = get_audio()

    msg = EmailMessage()
    msg["From"] = EMAIL_ADDRESS
    msg["To"] = receiver
    msg["Subject"] = "Message from JARVIS"
    msg.set_content(message_text)

    speak("Do you want to attach a file? Say yes or no.")
    if "yes" in get_audio():
        speak("Tell the file name.")
        filename = get_audio()
        if "." not in filename:
            filename += ".txt"

        file_path = find_file(filename)
        if file_path:
            with open(file_path, "rb") as f:
                msg.add_attachment(
                    f.read(),
                    maintype="application",
                    subtype="octet-stream",
                    filename=os.path.basename(file_path)
                )
            speak("File attached.")
        else:
            speak("File not found. Sending email without attachment.")

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
        server.send_message(msg)

    speak("Email sent successfully.")

#------------------- volume control ---------------------
def set_volume(level):
    """
    Set volume level (0 to 100)
    """
    devices = AudioUtilities.GetSpeakers()
    interface = devices.Activate(
        IAudioEndpointVolume._iid_,
        CLSCTX_ALL,
        None
    )

    volume = cast(interface, POINTER(IAudioEndpointVolume))

    # Convert 0–100 to 0.0–1.0
    volume.SetMasterVolumeLevelScalar(level / 100.0, None)

def volume_control(command):
    if "mute" in command:
        set_volume(0)
        speak("Volume muted")

    elif "increase volume" in command:
        set_volume(80)
        speak("Volume increased")

    elif "decrease volume" in command:
        set_volume(30)
        speak("Volume decreased")

    elif "set volume to" in command:
        try:
            level = int("".join(filter(str.isdigit, command)))
            level = max(0, min(level, 100))
            set_volume(level)
            speak(f"Volume set to {level} percent")
        except:
            speak("Unable to set volume")

# -------------------- Brightness ---------------------
def set_brightness(percent):
    sbc.set_brightness(percent)

def brightness_control(command):
    if "increase brightness" in command:
        set_brightness(80)
        speak("Brightness increased")

    elif "decrease brightness" in command:
        set_brightness(30)
        speak("Brightness decreased")

    elif "set brightness to" in command:
        try:
            level = int("".join(filter(str.isdigit, command)))
            level = max(0, min(level, 100))
            set_brightness(level)
            speak(f"Brightness set to {level} percent")
        except:
            speak("Unable to set brightness")

# ------------------- Whatsapp messenger ----------------
def open_whatsapp():
    service = Service(ChromeDriverManager().install())
    driver = webdriver.Chrome(service=service)
    driver.get("https://web.whatsapp.com")
    speak("Scan the QR code from your mobile")
    time.sleep(20)  # wait for QR scan
    return driver

def send_message(driver, name, message):
    try:
        search_box = driver.find_element(By.XPATH, '//div[@contenteditable="true"][@data-tab="3"]')
        search_box.click()
        search_box.send_keys(name)
        time.sleep(2)
        search_box.send_keys(Keys.ENTER)

        time.sleep(1)
        msg_box = driver.find_element(By.XPATH, '//div[@contenteditable="true"][@data-tab="10"]')
        msg_box.send_keys(message)
        msg_box.send_keys(Keys.ENTER)

        speak("Message sent successfully")

    except Exception as e:
        print(e)
        speak("Failed to send message")

# ------------------- Call --------------------
def open_phone_link():
    os.system("start ms-phone:")
    time.sleep(8)

def call_contact(name):
    if name not in CONTACTS:
        speak("Contact not found")
        return

    number = CONTACTS[name]
    speak(f"Calling {name}")

    open_phone_link()

    # Open dial pad
    pyautogui.hotkey("ctrl", "d")
    time.sleep(2)

    # Type number
    pyautogui.write(number, interval=0.1)
    time.sleep(1)

    # Press call
    pyautogui.press("enter")

def accept_call():
    speak("Accepting call")
    pyautogui.press("enter")

def end_call():
    speak("Ending call")
    pyautogui.hotkey("alt", "f4")

# ----------------- Main Program Logic ------------------
if __name__ == "__main__":    
    time_wish()
    speak("System activated.")

    while True:
        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            # Set a low energy threshold to detect the wake word more easily.
            recognizer.energy_threshold = 100

            try:
                # Listen for the wake word without showing anything on the console.
                print(f"Waiting for '{WAKE_WORD}'...")
                audio = recognizer.listen(source)
                command = recognizer.recognize_google(audio).lower()

                if WAKE_WORD in command:
                    print("Wake word detected!")
                    speak("Yes BOSS")

                    # Once the wake word is detected, get the next command.
                    user_command = get_audio()

                    if user_command == "timeout":
                        continue

                    if user_command:
                        print("You said:", user_command)
                        if "hello" in user_command:
                            speak("Hello there! How can I help you?")

                        elif "what is the time" in user_command:
                            speak(time.strftime("The time is %I:%M %p"))

                        elif "what is your name" in user_command:
                            speak("BOSS you can call me jarvis")

                        elif "tell me about you" in user_command:
                            speak("I am JARVIS, A personal AI assistant created by Shri Ram I can help you by controlling systems, analyzing data and responding intelligently through voice commands")
                        
                        elif "can you ask chat gpt" in user_command:
                            speak("what would you like to ask for chatgpt")
                            question = get_audio()
                            print("You:", question)
                            answer = chatgpt_search(question)
                            print("Jarvis:", answer)
                            speak(answer)

                        elif "mute" in user_command:
                            set_volume(0)

                        elif "increase volume" in user_command:
                            set_volume(80)

                        elif "decrease volume" in user_command:
                            set_volume(20)

                        elif "set volume to" in user_command:
                            try:
                                level = int("".join(filter(str.isdigit, command)))
                                level = max(0, min(level, 100))
                                set_volume(level)
                                speak(f"Volume set to {level} percent")
                            except:
                                speak("Unable to set volume")

                        elif "increase brightness" in user_command:
                            brightness_control(user_command)

                        elif "decrease brightness" in user_command:
                            brightness_control(user_command)

                        elif "set brightness to" in user_command:
                            brightness_control(user_command)

                        elif "open" in user_command:
                            os_control(user_command)

                        elif "search" in user_command:
                            browser_search(user_command)

                        elif "create file" in user_command:
                            create_file()

                        elif "find file" in user_command:
                            open_file()

                        elif "prepare ppt" in user_command:
                            ppt_creator()

                        elif "add slides" in user_command:
                            add_slide()

                        elif "convert file to pdf" in user_command or "convert to pdf" in user_command:
                            convert_to_pdf()

                        elif "can you send email" in user_command:
                            email_sender()

                        elif "send message to" in user_command and "saying" in user_command:
                            driver = open_whatsapp()

                            try:
                                # Remove starting phrase
                                user_command = user_command.replace("send message to", "").strip()

                                # Split name and message
                                name, message = user_command.split("saying", 1)

                                name = name.strip()
                                message = message.strip()

                                send_message(driver, name, message)

                            except Exception as e:
                                print(e)
                                speak("Please say the command clearly")

                        elif user_command.startswith("call"):
                            name = user_command.replace("call", "").strip()
                            call_contact(name)

                        elif "accept call" in user_command:
                            accept_call()

                        elif "end the call" in user_command or "cut call" in user_command:
                            end_call()

                        elif "good job jarvis" in user_command:
                            print("Thank you BOSS")
                            speak("Thank you BOSS")

                        elif "exit" in user_command or "bye" in user_command or "stop" in user_command:
                            speak("System Deactivated. Goodbye!")
                            break

                        elif user_command == "":
                            # Do nothing if no command was detected to avoid an endless loop
                            pass

                        else:
                            speak("sorry")


            except sr.UnknownValueError:
                # Silently ignore if wake word is not recognized.
                pass
            except sr.RequestError as e:
                print()

