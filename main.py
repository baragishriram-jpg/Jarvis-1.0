import speech_recognition as sr
import pyttsx3
import time
import datetime
import wikipedia
import os
import urllib.parse
import webbrowser
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import requests
from colorama import Fore,Style
from pptx.enum.text import PP_ALIGN
from reportlab.pdfgen import canvas
import win32com.client

# --- Configuration ---
WAKE_WORD = "jarvis"
LISTENING_DURATION = 10 # seconds 
LISTENING_TIMEOUT = 2 # seconds to wait for voice input
BASE_FILE_DIR = "C:/Users/admin/Desktop/Jarvis 3.0/Database"
BASE_PPT_DIR = "C:/Users/admin/Desktop/Jarvis 3.0/Database"
presentation = None
ppt_path = None

def print_loop():
    while True:
        print(Fore.BLUE +"listening...",end="",flush=True)
        print(Style.RESET_ALL,end="",flush=True)
        print("",end="",flush=True)


# --- Text-to-Speech Function ---
def speak(text):
    """
    Converts text to speech.

    Args:
        text (str): The text to be spoken.
    """
    engine = pyttsx3.init()
    engine.say(text)
    engine.runAndWait()


# --- Speech-to-Text Function ---
def get_audio():
    """
    Listens for a voice command and transcribes it to text.

    Returns:
        str or None: The transcribed text, or None if the speech was not recognized.
    """
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        # Calibrate for ambient noise.
        recognizer.adjust_for_ambient_noise(source)
        try:
            print("Listening...")
            audio = recognizer.listen(source, timeout=LISTENING_TIMEOUT)
            print("Recognizing...")
            text = recognizer.recognize_google(audio)
            return text.lower()
        except sr.UnknownValueError:
            return get_audio()
        except sr.WaitTimeoutError:
            return get_audio()
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")
            return None

# ---------------- OS CONTROL ----------------
def os_control(command):
    if "open chrome" in command:
        speak("Opening Google Chrome.")
        os.startfile("C:\\Program Files\\Google\\Chrome\\Application\\chrome.exe")

    elif "open edge" in command:
        speak("Opening Microsoft Edge.")
        os.startfile("C:\\Program Files (x86)\\Microsoft\\Edge\\Application\\msedge.exe")

    elif "open browser" in command:
        speak("Opening default browser.")
        webbrowser.open("https://www.google.com")

    elif "open youtube" in command:
        speak("Opening YouTube.")
        webbrowser.open("https://www.youtube.com")

    elif 'open google' in user_command:
        speak("Opening Google.")
        webbrowser.open("http://www.google.com")

    elif 'open instagram' in user_command:
        speak("Opening Instagram.")
        webbrowser.open("http://www.instagram.com")

    elif 'open github' in user_command:
        speak("Opening GitHub.")
        webbrowser.open("http://www.github.com")

    else:
        speak("OS command not recognized.")

# ---------------- BROWSER SEARCH ----------------
def browser_search(command):
    keywords = command.replace("search", "").replace("find", "").strip()
    if keywords:
        speak(f"Searching for {keywords}")
        query = urllib.parse.quote(keywords)
        webbrowser.open(f"https://www.google.com/search?q={query}")
    else:
        speak("What should I search for?")

# ------------- listener ---------------
def take_command(command):
    r = sr.Recognizer()
    with sr.Microphone() as source:
        r.adjust_for_ambient_noise(source, duration=2)
        print(command)
        speak(command)
        audio = r.listen(source)

    try:
        text = r.recognize_google(audio)
        print("You said:", text)
        return text.lower()
    except sr.UnknownValueError:
        speak("I didn't understand. Please say again.")
        return take_command(command)

# ---------- File Operations ----------
def create_file():
    filename = take_command("Tell the file name").replace(" ", "_")
    extension = take_command("Tell the file extension like txt, docx, pdf, python").replace(" ", "_")

    file_path = os.path.join(BASE_FILE_DIR, f"{filename}.{extension}")

    try:
        with open(file_path, "w") as f:
            f.write("File created by Jarvis voice command.")
        speak(f"{filename} dot {extension} created successfully")
    except Exception as e:
        speak("Unable to create the file")
        print(e)

def open_file():
    filename = take_command("Tell the file name").replace(" ", "_")
    extension = take_command("Tell the file extension").replace(" ", "_")

    file_path = os.path.join(BASE_FILE_DIR, f"{filename}.{extension}")

    if os.path.exists(file_path):
        os.startfile(file_path)
        speak(f"Opening {filename}")
    else:
        speak("File not found")

# ---------------- IMAGE FROM WEB ----------------
def get_image_from_web(topic):
    url = f"https://source.unsplash.com/800x600/?{topic}"
    response = requests.get(url, timeout=10)
    return BytesIO(response.content)

# ---------------- GET CONTENT FROM WEB ----------------
def get_content_from_web(topic):
    try:
        summary = wikipedia.summary(topic, sentences=4)
        return summary
    except:
        return f"Information about {topic} could not be fetched."

# ---------------- PPT CREATOR ----------------
def ppt_creator():
    global presentation, ppt_path

    topic = take_command("Please tell the presentation topic").lower()

    presentation = Presentation()
    file_name = f"{topic.replace(' ', '_')}.pptx"
    ppt_path = os.path.join(BASE_PPT_DIR, file_name)

    # TITLE SLIDE
    slide = presentation.slides.add_slide(presentation.slide_layouts[3])
    slide.shapes.title.text = topic.title()
    slide.placeholders[1].text = get_content_from_web(topic)
    slide.placeholders[1].pic = get_image_from_web(topic)

    speak(f"Presentation on {topic} created.")

    decision = take_command("Would you like to save this").lower().replace(" ","_")

    if "add_slide" in decision:
        add_slide(topic)
    elif "save_it" in decision:
        finish_presentation()

# ---------------- ADD SLIDE WITH AUTO CONTENT ----------------
def add_slide(topic):
    content_topic = take_command("Tell the slide topic").lower()

    slide = presentation.slides.add_slide(presentation.slide_layouts[3])
    slide.shapes.title.text = content_topic.title()

    slide.placeholders[1].pic = get_image_from_web(content_topic)
    slide.placeholders[1].text = get_content_from_web(content_topic)

    try:
        image = get_image_from_web(content_topic)
        slide.shapes.add_picture(
            image, Inches(4.5), Inches(1.5), width=Inches(4)
        )
    except:
        pass

    speak(f"Slide added for {content_topic}")

    decision = take_command("Do you want another slide or save").lower().replace(" ","_")

    if "add_another_slide" in decision:
        add_slide(topic)
    elif "save_it" in decision:
        finish_presentation()

# ---------------- SAVE PRESENTATION ----------------
def finish_presentation():
    global presentation, ppt_path

    if presentation is None:
        speak("Presentation is not created yet")
        return

    presentation.save(ppt_path)
    speak("Presentation saved successfully")

# ---------------- pdf converter -----------------
def convert_to_pdf():
    filename = take_command("Tell the file name").replace(" ", "_")
    ext = take_command("Tell the file extension").lower()

    input_path = os.path.join(BASE_FILE_DIR, f"{filename}.{ext}")
    output_path = os.path.join(BASE_FILE_DIR, f"{filename}.pdf")

    if not os.path.exists(input_path):
        speak("File not found")
        return

    # ---------- TXT TO PDF ----------
    if ext == "txt":
        try:
            c = canvas.Canvas(output_path)
            with open(input_path, "r", encoding="utf-8") as f:
                text = c.beginText(40, 800)
                for line in f:
                    text.textLine(line.strip())
                c.drawText(text)
            c.save()
            speak("Text file converted to PDF successfully")
        except:
            speak("Failed to convert text file")

    # ---------- WORD TO PDF ----------
    elif ext == "docx":
        try:
            word = win32com.client.Dispatch("Word.Application")
            doc = word.Documents.Open(input_path)
            doc.SaveAs(output_path, FileFormat=17)
            doc.Close()
            word.Quit()
            speak("Word file converted to PDF successfully")
        except:
            speak("Word conversion failed")

    # ---------- PPTX TO PDF ----------
    elif ext == "pptx":
        try:
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(input_path, WithWindow=False)
            presentation.SaveAs(output_path, 32)
            presentation.Close()
            powerpoint.Quit()
            speak("PowerPoint converted to PDF successfully")
        except:
            speak("PowerPoint conversion failed")

    else:
        speak("Unsupported file format")


# ----------------- Main Program Logic ------------------
if __name__ == "__main__":
    speak("System activated Sir.")

    while True:
        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            # Set a low energy threshold to detect the wake word more easily.
            recognizer.energy_threshold = 4000

            try:
                # Listen for the wake word without showing anything on the console.
                print(f"Waiting for '{WAKE_WORD}'...")
                audio = recognizer.listen(source)
                command = recognizer.recognize_google(audio).lower()

                if WAKE_WORD in command:
                    print("Wake word detected!")
                    speak("Yes Sir?")

                    # Once the wake word is detected, get the next command.
                    user_command = get_audio()

                    if user_command == "timeout":
                        continue

                    if user_command:
                        print("You said:", user_command)
                        if "hello" in user_command:
                            speak("Hello there! How can I help you?")
                        elif "stop listening" in user_command:
                            speak("Goodbye!")
                            break
                        if 'wikipedia' in user_command:
                            speak("Searching Wikipedia...")
                            command = command.replace("wikipedia", "")
                            try:
                                # Use a try-except block to handle potential errors
                                results = wikipedia.summary(command, sentences=2)
                                speak("According to Wikipedia...")
                                speak(results)
                            except wikipedia.exceptions.PageError:
                                speak("Sorry, I couldn't find a page for that on Wikipedia.")
                            except wikipedia.exceptions.DisambiguationError as e:
                                speak("Please be more specific. Here are a few options:")
                                speak(str(e.options[:5]))  # Speak the first 5 options.

                        elif "open" in user_command:
                            os_control(user_command)

                        elif "search" in user_command:
                            browser_search(user_command)

                        elif "create file" in user_command:
                            create_file()

                        elif "find file" in user_command:
                            open_file()

                        elif "prepare ppt" in user_command:
                            ppt_creator()

                        elif "add slides" in user_command:
                            add_slide()

                        elif "convert file to pdf" in user_command or "convert to pdf" in user_command:
                            convert_to_pdf()

                        elif 'the time' in user_command:
                            str_time = datetime.datetime.now().strftime("%I:%M %p")
                            speak(f"The current time is {str_time}")

                        elif 'what is your name' in user_command:
                            speak("you can call me jarvis")


                        elif 'exit' in user_command or 'bye' in user_command or 'stop' in user_command:
                            speak("Stopping system. Goodbye!")
                            break

                        elif user_command == "":
                            # Do nothing if no command was detected to avoid an endless loop
                            pass

                        else:
                            speak("sorry")


            except sr.UnknownValueError:
                # Silently ignore if wake word is not recognized.
                pass
            except sr.RequestError as e:
                print()

